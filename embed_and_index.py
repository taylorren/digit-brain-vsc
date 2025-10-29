import json
import os
from tqdm import tqdm
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss

import fitz  # PyMuPDF
from pptx import Presentation

# Helper to extract text from PPTX
def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PPTX {pptx_path}: {e}")
        return ""

# Helper to extract text from PDF
def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    except Exception as e:
        print(f"Error reading PDF {pdf_path}: {e}")
        return ""


# === 配置区 ===
# 指定要扫描的多个根目录（直接在此处硬编码路径，按需修改）
SCAN_ROOTS = [
    r"F:/My Books/Working/My Own Writings Managed by Obsidian",  # 主目录
    r"F:/My Books/Working/_各读书会",  # 示例：添加更多目录
    # r"D:/Some/Other/Path"        # 示例：再加一个目录
]

# === 扫描文件 ===
md_files = []
pdf_files = []
pptx_files = []
for scan_root in SCAN_ROOTS:
    for root, dirs, files in os.walk(scan_root):
        for file in files:
            if file.endswith('.md'):
                md_files.append(os.path.join(root, file))
            elif file.endswith('.pdf'):
                pdf_files.append(os.path.join(root, file))
            elif file.endswith('.pptx'):
                pptx_files.append(os.path.join(root, file))

# Unified progress indicator for all file types
md_data = []
pdf_data = []
pptx_data = []
total_md = len(md_files)
total_pdf = len(pdf_files)
total_pptx = len(pptx_files)
md_done = pdf_done = pptx_done = 0
def print_progress():
    print(f"processing {md_done}/{total_md} md files, {pdf_done}/{total_pdf} pdf files, {pptx_done}/{total_pptx} pptx files", end='\r')
print_progress()
for path in md_files:
    try:
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        content_stripped = content.strip()
        if len(content_stripped) >= 300:
            md_data.append({'path': path, 'content': content, 'type': 'md'})
        else:
            print(f"[跳过过短md] {path} ({len(content_stripped)} chars)")
    except Exception as e:
        print(f"Error reading {path}: {e}")
    md_done += 1
    print_progress()
for path in pdf_files:
    content = extract_text_from_pdf(path)
    if content.strip():
        pdf_data.append({'path': path, 'content': content, 'type': 'pdf'})
    pdf_done += 1
    print_progress()
for path in pptx_files:
    content = extract_text_from_pptx(path)
    if content.strip():
        pptx_data.append({'path': path, 'content': content, 'type': 'pptx'})
    pptx_done += 1
    print_progress()
print()  # Newline after progress

# Combine all data
all_data = md_data + pdf_data + pptx_data

# --- Model selection: use local bge-large-zh model ---
model = SentenceTransformer(os.path.join(os.path.dirname(__file__), 'models', 'bge-large-zh'))

# Generate embeddings for each document (using the content field) with progress indicator
embeddings = []
for i, doc in enumerate(tqdm(all_data, desc='Embedding documents', unit='doc')):
    emb = model.encode(doc['content'], show_progress_bar=True, normalize_embeddings=True)  # normalize for cosine
    embeddings.append(emb)
    if (i + 1) % 10 == 0 or (i + 1) == len(all_data):
        print(f"Embedded {i + 1}/{len(all_data)} documents", end='\r')
embeddings = np.vstack(embeddings)

# --- Use FAISS IndexFlatIP for cosine similarity ---
dimension = embeddings.shape[1]
index = faiss.IndexFlatIP(dimension)  # Inner Product = Cosine if normalized
index.add(embeddings)

# Save FAISS index and metadata
faiss.write_index(index, os.path.join(os.path.dirname(__file__), 'md_faiss.index'))
with open(os.path.join(os.path.dirname(__file__), 'md_faiss_meta.json'), 'w', encoding='utf-8') as f:
    json.dump(all_data, f, ensure_ascii=False, indent=2)

print(f"Indexed {len(all_data)} documents (md + pdf + pptx) with dimension {dimension}.")
print("FAISS index and metadata saved.")
